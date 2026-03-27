#!/usr/bin/env python3
"""
PII Shield - Local PII/PHI Anonymizer for Safe LLM Processing
============================================================
Strips sensitive info from documents before sending to Claude.
Stores a local mapping so you can restore originals later.

Usage:
  python pii_shield.py anonymize report.pdf
  python pii_shield.py anonymize spreadsheet.xlsx --output clean.txt
  python pii_shield.py restore clean.txt --session abc123
  python pii_shield.py restore report.docx --session abc123 --output restored.docx
  python pii_shield.py sessions
"""

import sys
import os
import json
import re
import hashlib
import uuid
import argparse
from datetime import datetime
from pathlib import Path
from typing import Optional

# ─── Lazy imports (checked at runtime with helpful errors) ───────────────────

def require(pkg, install_name=None):
    """Import a package with a friendly error if missing."""
    import importlib
    try:
        return importlib.import_module(pkg)
    except ImportError:
        name = install_name or pkg
        print(f"\n❌  Missing package: {pkg}")
        print(f"    Run:  pip install {name}")
        sys.exit(1)


# ─── Constants ───────────────────────────────────────────────────────────────

MAPPINGS_DIR = Path.home() / ".pii_shield" / "sessions"
MAPPINGS_DIR.mkdir(parents=True, exist_ok=True)

# Entity types to detect (Presidio built-ins + custom PHI)
DEFAULT_ENTITIES = [
    "PERSON",
    "EMAIL_ADDRESS",
    "PHONE_NUMBER",
    "US_SSN",
    "US_PASSPORT",
    "US_DRIVER_LICENSE",
    "CREDIT_CARD",
    "IBAN_CODE",
    "IP_ADDRESS",
    "URL",
    "US_BANK_NUMBER",
    "MEDICAL_LICENSE",
    "DATE_TIME",
    "LOCATION",
    "NRP",               # Nationality, Religion, Political group
    "ORGANIZATION",
    "US_ITIN",
]

# Labels shown in the anonymized output
ENTITY_LABELS = {
    "PERSON": "PERSON",
    "EMAIL_ADDRESS": "EMAIL",
    "PHONE_NUMBER": "PHONE",
    "US_SSN": "SSN",
    "US_PASSPORT": "PASSPORT",
    "US_DRIVER_LICENSE": "DL",
    "CREDIT_CARD": "CC",
    "IBAN_CODE": "IBAN",
    "IP_ADDRESS": "IP",
    "URL": "URL",
    "US_BANK_NUMBER": "BANK_ACCT",
    "MEDICAL_LICENSE": "MED_LIC",
    "DATE_TIME": "DATE",
    "LOCATION": "LOCATION",
    "NRP": "NRP",
    "ORGANIZATION": "ORG",
    "US_ITIN": "ITIN",
}


# ─── Document Parsers ─────────────────────────────────────────────────────────

def parse_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def parse_pdf(path: Path) -> str:
    pdfplumber = require("pdfplumber")
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def parse_docx(path: Path) -> str:
    docx = require("docx", "python-docx")
    doc = docx.Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also grab tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def parse_excel(path: Path) -> str:
    pd = require("pandas")
    require("openpyxl")
    xl = pd.ExcelFile(str(path))
    parts = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        parts.append(f"[Sheet: {sheet}]")
        parts.append(df.to_string(index=False))
    return "\n\n".join(parts)


def parse_csv(path: Path) -> str:
    pd = require("pandas")
    df = pd.read_csv(str(path))
    return df.to_string(index=False)


def parse_pptx(path: Path) -> str:
    pptx = require("pptx", "python-pptx")
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def parse_document(path: Path) -> str:
    """Dispatch to the right parser based on file extension."""
    ext = path.suffix.lower()
    parsers = {
        ".txt":  parse_txt,
        ".md":   parse_txt,
        ".pdf":  parse_pdf,
        ".docx": parse_docx,
        ".doc":  parse_docx,
        ".xlsx": parse_excel,
        ".xls":  parse_excel,
        ".xlsm": parse_excel,
        ".csv":  parse_csv,
        ".pptx": parse_pptx,
        ".ppt":  parse_pptx,
    }
    if ext not in parsers:
        print(f"⚠️  Unsupported file type: {ext}")
        print(f"   Supported: {', '.join(parsers)}")
        sys.exit(1)
    print(f"📄  Parsing {path.name} …")
    return parsers[ext](path)


# ─── PII Engine ───────────────────────────────────────────────────────────────

def build_analyzer(entities=None):
    """Build a Presidio AnalyzerEngine with spaCy NLP backend."""
    from presidio_analyzer import AnalyzerEngine
    from presidio_analyzer.nlp_engine import NlpEngineProvider

    # Try large model first, fall back to small
    for model in ("en_core_web_lg", "en_core_web_sm", "en_core_web_md"):
        try:
            import spacy
            spacy.load(model)
            nlp_config = {"nlp_engine_name": "spacy", "models": [{"lang_code": "en", "model_name": model}]}
            print(f"🧠  Using spaCy model: {model}")
            break
        except OSError:
            continue
    else:
        print("❌  No spaCy model found. Run one of:")
        print("    python -m spacy download en_core_web_lg")
        print("    python -m spacy download en_core_web_sm")
        sys.exit(1)

    provider = NlpEngineProvider(nlp_configuration=nlp_config)
    nlp_engine = provider.create_engine()
    return AnalyzerEngine(nlp_engine=nlp_engine, supported_languages=["en"])


def anonymize_text(text: str, entities=None, score_threshold=0.4) -> tuple[str, dict]:
    """
    Detect and replace PII/PHI in text.
    Returns (anonymized_text, mapping_dict).
    """
    from presidio_analyzer import AnalyzerEngine
    from presidio_anonymizer import AnonymizerEngine
    from presidio_anonymizer.entities import OperatorConfig

    target_entities = entities or DEFAULT_ENTITIES

    print("🔍  Analyzing for PII/PHI …")
    analyzer = build_analyzer()
    anonymizer = AnonymizerEngine()

    results = analyzer.analyze(
        text=text,
        entities=target_entities,
        language="en",
        score_threshold=score_threshold,
    )

    if not results:
        print("✅  No PII/PHI detected.")
        return text, {}

    # Sort by position (reverse) so we can replace without offset issues
    results = sorted(results, key=lambda r: r.start, reverse=True)

    # Build mapping: original_value → placeholder
    mapping = {}          # placeholder → original
    counters = {}         # entity_type → count
    anonymized = text

    for result in results:
        original = text[result.start:result.end]

        # Deduplicate: same value → same placeholder
        key = (result.entity_type, original.strip().lower())
        if key not in counters:
            label = ENTITY_LABELS.get(result.entity_type, result.entity_type)
            idx = len([v for v in mapping.values() if v["entity_type"] == result.entity_type]) + 1
            placeholder = f"[{label}_{idx}]"
            counters[key] = placeholder
            mapping[placeholder] = {
                "original": original,
                "entity_type": result.entity_type,
                "score": round(result.score, 3),
            }
        else:
            placeholder = counters[key]

        anonymized = anonymized[:result.start] + placeholder + anonymized[result.end:]

    return anonymized, mapping


# ─── Session Storage ──────────────────────────────────────────────────────────

def save_session(mapping: dict, source_file: str, session_id: str = None) -> str:
    """Persist a mapping to disk under a session ID."""
    sid = session_id or str(uuid.uuid4())[:8]
    session = {
        "session_id": sid,
        "created": datetime.now().isoformat(),
        "source_file": source_file,
        "entity_count": len(mapping),
        "mapping": mapping,
    }
    out = MAPPINGS_DIR / f"{sid}.json"
    out.write_text(json.dumps(session, indent=2))
    return sid, out


def load_session(session_id: str) -> dict:
    """Load a session from disk."""
    path = MAPPINGS_DIR / f"{session_id}.json"
    if not path.exists():
        print(f"❌  Session not found: {session_id}")
        print(f"   Sessions are stored in: {MAPPINGS_DIR}")
        sys.exit(1)
    return json.loads(path.read_text())


def list_sessions():
    """Print all stored sessions."""
    files = sorted(MAPPINGS_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not files:
        print("No sessions stored yet.")
        return
    print(f"\n{'SESSION ID':<12} {'CREATED':<22} {'ENTITIES':<10} SOURCE FILE")
    print("─" * 70)
    for f in files:
        try:
            s = json.loads(f.read_text())
            created = s.get("created", "")[:19].replace("T", " ")
            print(f"{s['session_id']:<12} {created:<22} {s['entity_count']:<10} {s.get('source_file','?')}")
        except Exception:
            pass


# ─── Restore ─────────────────────────────────────────────────────────────────

def restore_text(anonymized: str, mapping: dict) -> str:
    """Replace placeholders back with original values (plain string)."""
    result = anonymized
    for placeholder, info in mapping.items():
        result = result.replace(placeholder, info["original"])
    return result


def _replace_in_str(value: str, mapping: dict) -> str:
    """Replace all known placeholders inside a single string value."""
    if not isinstance(value, str):
        return value
    for placeholder, info in mapping.items():
        value = value.replace(placeholder, info["original"])
    return value


# ── Per-format restore writers ────────────────────────────────────────────────

def restore_txt(path: Path, mapping: dict, out_path: Path):
    """Restore a plain-text or markdown file."""
    text = path.read_text(encoding="utf-8", errors="replace")
    out_path.write_text(restore_text(text, mapping), encoding="utf-8")


def restore_docx(path: Path, mapping: dict, out_path: Path):
    """
    Restore a .docx by walking every run in every paragraph and table cell.
    Preserves all formatting (bold, font, size, colour, etc.) because we
    operate run-by-run rather than on the raw XML string.
    """
    docx_mod = require("docx", "python-docx")
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))

    def fix_paragraph(para):
        # First pass: collect the full paragraph text so we can detect
        # placeholders that were split across multiple runs by Word's XML
        # (e.g. "[PERSON" in one run and "_1]" in the next).
        full = para.text
        if not any(ph in full for ph in mapping):
            return  # nothing to do

        # Simple case: placeholder is wholly inside a single run
        for run in para.runs:
            run.text = _replace_in_str(run.text, mapping)

        # Harder case: placeholder was split across runs.
        # Rebuild by merging all run text into run[0], clearing the rest.
        # We only do this when the paragraph still contains a bracket after
        # the per-run pass (meaning a split placeholder survived).
        if "[" in para.text:
            combined = "".join(r.text for r in para.runs)
            fixed = _replace_in_str(combined, mapping)
            if fixed != combined and para.runs:
                para.runs[0].text = fixed
                for run in para.runs[1:]:
                    run.text = ""

    # Walk all paragraphs in the body
    for para in doc.paragraphs:
        fix_paragraph(para)

    # Walk all tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    fix_paragraph(para)

    # Walk headers and footers
    for section in doc.sections:
        for hf in (section.header, section.footer,
                   section.even_page_header, section.even_page_footer,
                   section.first_page_header, section.first_page_footer):
            if hf is not None:
                for para in hf.paragraphs:
                    fix_paragraph(para)
                for table in hf.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                fix_paragraph(para)

    doc.save(str(out_path))


def restore_xlsx(path: Path, mapping: dict, out_path: Path):
    """
    Restore an .xlsx by iterating every cell in every sheet.
    Preserves cell types — only string cells are touched.
    """
    openpyxl = require("openpyxl")
    wb = openpyxl.load_workbook(str(path))
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell.value = _replace_in_str(cell.value, mapping)
    wb.save(str(out_path))


def restore_csv(path: Path, mapping: dict, out_path: Path):
    """Restore a CSV — treat it as text so we don't need pandas."""
    text = path.read_text(encoding="utf-8", errors="replace")
    out_path.write_text(restore_text(text, mapping), encoding="utf-8")


def restore_pptx(path: Path, mapping: dict, out_path: Path):
    """
    Restore a .pptx by walking every text run on every slide.
    Also handles slide notes and slide masters.
    """
    pptx_mod = require("pptx", "python-pptx")
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(path))

    def fix_text_frame(tf):
        for para in tf.paragraphs:
            # Per-run pass
            for run in para.runs:
                run.text = _replace_in_str(run.text, mapping)
            # Split-run fallback
            if "[" in para.text:
                combined = "".join(r.text for r in para.runs)
                fixed = _replace_in_str(combined, mapping)
                if fixed != combined and para.runs:
                    para.runs[0].text = fixed
                    for run in para.runs[1:]:
                        run.text = ""

    def fix_shape(shape):
        if shape.has_text_frame:
            fix_text_frame(shape.text_frame)
        # Groups of shapes
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for s in shape.shapes:
                fix_shape(s)

    for slide in prs.slides:
        for shape in slide.shapes:
            fix_shape(shape)
        # Slide notes
        if slide.has_notes_slide:
            fix_text_frame(slide.notes_slide.notes_text_frame)

    prs.save(str(out_path))


def restore_pdf(path: Path, mapping: dict, out_path: Path):
    """
    PDFs are binary and layout-dependent — true in-place restoration isn't
    possible without recreating the file. We restore to a UTF-8 .txt file
    (or a .pdf if reportlab is available) with the original values filled in.
    """
    # Extract text
    pdfplumber = require("pdfplumber")
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            pages.append(f"[Page {i}]\n{text}")
    raw = "\n\n".join(pages)
    restored = restore_text(raw, mapping)

    # Try to write a proper PDF via reportlab; fall back to .txt
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph as RLPara, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch

        # Force .pdf extension on output
        pdf_out = out_path.with_suffix(".pdf")
        doc_rl = SimpleDocTemplate(str(pdf_out), pagesize=letter,
                                   leftMargin=inch, rightMargin=inch,
                                   topMargin=inch, bottomMargin=inch)
        styles = getSampleStyleSheet()
        story = []
        for line in restored.splitlines():
            story.append(RLPara(line or "&nbsp;", styles["Normal"]))
            story.append(Spacer(1, 4))
        doc_rl.build(story)
        return pdf_out  # Caller uses this path in the success message

    except ImportError:
        # reportlab not installed — write plain text
        txt_out = out_path.with_suffix(".txt")
        txt_out.write_text(restored, encoding="utf-8")
        print("ℹ️   reportlab not found; restored PDF saved as plain text.")
        print("    Install reportlab for PDF output:  pip install reportlab")
        return txt_out


# ── Dispatcher ────────────────────────────────────────────────────────────────

def restore_file(path: Path, mapping: dict, out_path: Path) -> Path:
    """
    Route to the right restore function based on file extension.
    Returns the actual output path (may differ for PDFs).

    Supported formats
    ─────────────────
    Text   : .txt  .md  .log  .json  .xml  .html  .htm  .yaml  .yml  .toml  .ini  .cfg
    Word   : .docx  .doc
    Excel  : .xlsx  .xlsm  .xls
    CSV    : .csv  .tsv
    PowerPoint: .pptx  .ppt
    PDF    : .pdf  (restored as PDF if reportlab installed, else .txt)
    """
    ext = path.suffix.lower()

    # Plain-text family (all read/writable as UTF-8 strings)
    TEXT_EXTS = {".txt", ".md", ".log", ".json", ".xml",
                 ".html", ".htm", ".yaml", ".yml",
                 ".toml", ".ini", ".cfg", ".rst", ".rtf"}

    if ext in TEXT_EXTS:
        restore_txt(path, mapping, out_path)

    elif ext in {".docx", ".doc"}:
        restore_docx(path, mapping, out_path)

    elif ext in {".xlsx", ".xlsm", ".xls"}:
        restore_xlsx(path, mapping, out_path)

    elif ext in {".csv", ".tsv"}:
        restore_csv(path, mapping, out_path)

    elif ext in {".pptx", ".ppt"}:
        restore_pptx(path, mapping, out_path)

    elif ext == ".pdf":
        actual = restore_pdf(path, mapping, out_path)
        return actual  # may have changed extension

    else:
        print(f"⚠️  No binary-aware restore available for '{ext}'.")
        print("    Attempting plain-text fallback (may corrupt binary files).")
        restore_txt(path, mapping, out_path)

    return out_path


# ─── CLI Commands ─────────────────────────────────────────────────────────────

def cmd_anonymize(args):
    path = Path(args.file)
    if not path.exists():
        print(f"❌  File not found: {path}")
        sys.exit(1)

    # Parse document
    text = parse_document(path)
    print(f"📝  Extracted {len(text):,} characters")

    # Anonymize
    anonymized, mapping = anonymize_text(
        text,
        score_threshold=args.threshold,
    )

    # Save session
    sid, session_path = save_session(mapping, str(path))

    # Determine output path
    if args.output:
        out_path = Path(args.output)
    else:
        out_path = path.parent / f"{path.stem}_anonymized.txt"

    out_path.write_text(anonymized, encoding="utf-8")

    # Summary
    print(f"\n{'─'*55}")
    print(f"✅  Anonymization complete!")
    print(f"{'─'*55}")
    print(f"  Session ID  : {sid}  (save this to restore later)")
    print(f"  Output file : {out_path}")
    print(f"  Mapping file: {session_path}")
    print(f"  Entities found: {len(mapping)}")
    if mapping:
        print(f"\n  Replacements made:")
        for ph, info in sorted(mapping.items()):
            orig = info['original']
            if len(orig) > 40:
                orig = orig[:37] + "…"
            print(f"    {ph:<20} ← \"{orig}\"  (conf: {info['score']})")
    print(f"\n  📋 The anonymized file is ready to paste into Claude.")
    print(f"  🔑 Keep your session ID to restore: {sid}")
    print()


def cmd_restore(args):
    # Load session
    session = load_session(args.session)
    mapping = session["mapping"]

    # ── Stdin path (no file given) ────────────────────────────────────────────
    if not args.file:
        print("Paste the anonymized text below. Press Ctrl+D (or Ctrl+Z on Windows) when done:\n")
        anonymized = sys.stdin.read()
        restored = restore_text(anonymized, mapping)
        if args.output:
            out_path = Path(args.output)
            out_path.write_text(restored, encoding="utf-8")
            print(f"✅  Restored text written to: {out_path}")
        else:
            print("\n" + "─"*55)
            print("RESTORED TEXT:")
            print("─"*55)
            print(restored)
        return

    # ── File path ─────────────────────────────────────────────────────────────
    path = Path(args.file)
    if not path.exists():
        print(f"❌  File not found: {path}")
        sys.exit(1)

    # Determine output path
    if args.output:
        out_path = Path(args.output)
    else:
        out_path = path.parent / f"{path.stem}_restored{path.suffix}"

    print(f"📄  Restoring {path.name} …")
    actual_out = restore_file(path, mapping, out_path)

    replacements = sum(
        1 for info in mapping.values()
        if info["original"]  # non-empty originals
    )
    print(f"\n{'─'*55}")
    print(f"✅  Restore complete!")
    print(f"{'─'*55}")
    print(f"  Input      : {path}")
    print(f"  Output     : {actual_out}")
    print(f"  Placeholders replaced: {replacements}")
    print()


def cmd_sessions(args):
    print(f"📁 Sessions stored in: {MAPPINGS_DIR}")
    list_sessions()


def cmd_inspect(args):
    """Show what's in a session mapping."""
    session = load_session(args.session)
    print(f"\nSession: {session['session_id']}")
    print(f"Created: {session['created']}")
    print(f"Source:  {session['source_file']}")
    print(f"\n{'PLACEHOLDER':<22} {'TYPE':<22} {'CONF':<6} ORIGINAL VALUE")
    print("─" * 80)
    for ph, info in sorted(session["mapping"].items()):
        orig = info['original'].replace('\n', ' ')
        if len(orig) > 35:
            orig = orig[:32] + "…"
        print(f"{ph:<22} {info['entity_type']:<22} {info['score']:<6} {orig}")


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        prog="pii_shield",
        description="🛡️  PII Shield – Anonymize documents before sending to Claude",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pii_shield.py anonymize patient_report.pdf
  python pii_shield.py anonymize data.xlsx --output clean.txt --threshold 0.5
  python pii_shield.py restore clean.txt --session abc12345
  python pii_shield.py restore report.docx --session abc12345
  python pii_shield.py restore report.docx --session abc12345 --output report_real.docx
  python pii_shield.py restore slides.pptx --session abc12345
  python pii_shield.py restore data.xlsx --session abc12345
  python pii_shield.py sessions
  python pii_shield.py inspect --session abc12345
        """,
    )
    sub = parser.add_subparsers(dest="command", required=True)

    # anonymize
    p_anon = sub.add_parser("anonymize", help="Strip PII/PHI from a document")
    p_anon.add_argument("file", help="Path to document (pdf, docx, xlsx, csv, pptx, txt, …)")
    p_anon.add_argument("-o", "--output", help="Output file path (default: <name>_anonymized.txt)")
    p_anon.add_argument("--threshold", type=float, default=0.4,
                        help="Confidence threshold 0-1 (default: 0.4, lower = more aggressive)")
    p_anon.set_defaults(func=cmd_anonymize)

    # restore
    p_rest = sub.add_parser("restore", help="Restore original values into a file")
    p_rest.add_argument("file", nargs="?",
                        help="Anonymized file — any supported format (omit to paste from stdin)")
    p_rest.add_argument("-s", "--session", required=True, help="Session ID from anonymization")
    p_rest.add_argument("-o", "--output",
                        help="Output file (default: <name>_restored.<ext>)")
    p_rest.set_defaults(func=cmd_restore)

    # sessions
    p_sess = sub.add_parser("sessions", help="List all stored anonymization sessions")
    p_sess.set_defaults(func=cmd_sessions)

    # inspect
    p_insp = sub.add_parser("inspect", help="Show the mapping table for a session")
    p_insp.add_argument("-s", "--session", required=True, help="Session ID to inspect")
    p_insp.set_defaults(func=cmd_inspect)

    args = parser.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
